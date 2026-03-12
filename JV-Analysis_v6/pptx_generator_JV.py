"""
PowerPoint slide generator for JV Analysis.

Produces a single 16:9 slide matching the KIT lab template:
  - Title placeholder:        "Process name"
  - Stack line:               "Stack:  ..."
  - Left image (mid section): JV Curves -- Best Device per Condition
  - Right image (mid section):Combined Boxplot
  - Bottom-left table:        V_OC / J_SC / FF / PCE  x  Baseline | Current batch
  - Bottom-right text area:   "Speculations:"
  - Footer:                   "Batch ID: ..."

Images are placed with a fixed left/top anchor and a fixed HEIGHT only.
Width is derived automatically from each PNG's natural aspect ratio so
plots are never squeezed or stretched.
"""

import io


def _fig_to_png_bytes(fig):
    """
    Export a Plotly figure to a high-resolution PNG.

    We deliberately export at a large fixed pixel size (1400 x 860) rather
    than reading fig.layout.width/height.  The stored layout dimensions are
    the values set at figure-creation time and do NOT reflect how large the
    user has stretched the plot in the browser.  Exporting at a generous size
    keeps axis labels, tick text and legend text sharp and readable on the
    slide regardless of the original figure dimensions.

    width=None is used in add_picture so python-pptx derives the slide width
    from the PNG's natural aspect ratio -- nothing is ever squished.
    """
    import plotly.io as pio

    return pio.to_image(fig, format="png", width=1400, height=860, scale=2)


def generate_jv_pptx_bytes(jv_fig, boxplot_fig, batch_ids):
    """
    Build a single-slide PowerPoint and return the file as bytes.

    Parameters
    ----------
    jv_fig : plotly.graph_objects.Figure or None
    boxplot_fig : plotly.graph_objects.Figure or None
    batch_ids : list[str]

    Returns
    -------
    bytes
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    # ── Colour palette ────────────────────────────────────────────────────────
    NAVY       = RGBColor(0x1F, 0x39, 0x64)
    TEAL_DARK  = RGBColor(0x1B, 0x6B, 0x5E)
    TEAL_MID   = RGBColor(0xD5, 0xEA, 0xE6)
    TEAL_LIGHT = RGBColor(0xEB, 0xF4, 0xF2)
    DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

    # ── Presentation setup ────────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # ── Helper: plain text box ────────────────────────────────────────────────
    def _textbox(left, top, w, h, text, size, bold=False,
                 color=DARK_GRAY, align=PP_ALIGN.LEFT, italic=False):
        tb = slide.shapes.add_textbox(left, top, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text           = text
        run.font.size      = Pt(size)
        run.font.bold      = bold
        run.font.italic    = italic
        run.font.color.rgb = color
        return tb

    # ── Layout constants ──────────────────────────────────────────────────────
    M   = Inches(0.28)   # left margin
    MID = Inches(6.80)   # x-start of right column
    RW  = Inches(6.25)   # right-column content width

    # ── Title ─────────────────────────────────────────────────────────────────
    _textbox(M, Inches(0.10), Inches(12.77), Inches(0.70),
             "Process name", size=28, bold=True, color=NAVY)

    # ── Stack label ───────────────────────────────────────────────────────────
    _textbox(M, Inches(0.83), Inches(12.77), Inches(0.40),
             "Stack:  ...", size=14, color=DARK_GRAY)

    # ── Figures ───────────────────────────────────────────────────────────────
    # Only height is specified; width=None lets python-pptx preserve the PNG's
    # natural aspect ratio automatically.
    IMG_TOP = Inches(1.32)
    IMG_H   = Inches(3.45)

    def _add_figure(fig, left_x):
        if fig is None:
            return
        try:
            png = _fig_to_png_bytes(fig)
            slide.shapes.add_picture(io.BytesIO(png), left_x, IMG_TOP,
                                     width=None, height=IMG_H)
        except Exception:
            pass

    _add_figure(jv_fig,      M)    # left  – aspect ratio from its own PNG
    _add_figure(boxplot_fig, MID)  # right – aspect ratio from its own PNG

    # ── Performance table (bottom left) ──────────────────────────────────────
    TBL_LEFT = M
    TBL_TOP  = Inches(4.95)
    TBL_W    = Inches(6.0)
    TBL_H    = Inches(2.00)
    N_ROWS   = 5   # 1 header + 4 data rows
    N_COLS   = 3

    row_labels  = ["V_OC", "J_SC", "FF", "PCE"]
    col_headers = ["", "Baseline", "Current batch"]

    tbl = slide.shapes.add_table(N_ROWS, N_COLS,
                                 TBL_LEFT, TBL_TOP, TBL_W, TBL_H).table
    tbl.columns[0].width = Inches(1.20)
    tbl.columns[1].width = Inches(2.40)
    tbl.columns[2].width = Inches(2.40)

    def _cell_fill(cell, rgb):
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        sf   = etree.SubElement(tcPr, qn("a:solidFill"))
        sc   = etree.SubElement(sf,   qn("a:srgbClr"))
        sc.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")

    def _cell_text(cell, text, size, bold=False,
                   color=DARK_GRAY, align=PP_ALIGN.CENTER):
        cell.text = ""
        tf  = cell.text_frame
        tf.word_wrap = False
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text           = text
        run.font.size      = Pt(size)
        run.font.bold      = bold
        run.font.color.rgb = color

    # Header row
    for ci, hdr in enumerate(col_headers):
        cell = tbl.cell(0, ci)
        _cell_fill(cell, TEAL_DARK)
        _cell_text(cell, hdr, size=12, bold=True, color=WHITE)

    # Data rows – alternating teal shades
    row_fills = [TEAL_MID, TEAL_LIGHT, TEAL_MID, TEAL_LIGHT]
    for ri, label in enumerate(row_labels):
        for ci in range(N_COLS):
            cell = tbl.cell(ri + 1, ci)
            _cell_fill(cell, row_fills[ri])
            _cell_text(cell, label if ci == 0 else "",
                       size=12, color=DARK_GRAY)

    # ── Speculations text area (bottom right) ─────────────────────────────────
    _textbox(MID, Inches(4.95), RW, Inches(2.30),
             "Speculations:", size=16, color=TEAL_DARK)

    # ── Batch ID footer ───────────────────────────────────────────────────────
    batch_str = ", ".join(str(b) for b in batch_ids) if batch_ids else "-"
    _textbox(M, Inches(7.20), Inches(12.77), Inches(0.28),
             f"Batch ID:   {batch_str}", size=11, bold=True, color=DARK_GRAY)

    # ── Serialize ─────────────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
